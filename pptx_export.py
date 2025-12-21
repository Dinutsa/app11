import io
import os
import textwrap
import pandas as pd  
import matplotlib
matplotlib.use('Agg') 
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn
from classification import QuestionInfo, QuestionType
from summary import QuestionSummary
from typing import List

LOGO_FILE = "logo2.png"
UNIV_NAME = "Чернівецький національний університет імені Юрія Фкдьковича"
CHART_DPI = 150
FONT_SIZE_CHART = 11        
FONT_SIZE_HEADER = 12 
FONT_SIZE_DATA = 11   
BAR_WIDTH = 0.6

def set_table_grid_style(table):
    tbl = table._tbl
    tblPr = tbl.tblPr
    tblStyle = tblPr.find(qn('a:tableStyleId'))
    if tblStyle is None:
        tblStyle = OxmlElement('a:tableStyleId')
        tblPr.append(tblStyle)
    tblStyle.text = '{5940675A-B579-460E-94D1-54222C63F5DA}'

def create_chart_image(qs: QuestionSummary) -> io.BytesIO:
    plt.close('all') 
    plt.clf()
    plt.rcParams.update({'font.size': FONT_SIZE_CHART})
    
    labels = qs.table["Варіант відповіді"].astype(str).tolist()
    values = qs.table["Кількість"]
    wrapped_labels = [textwrap.fill(l, 25) for l in labels]

    is_scale = (qs.question.qtype == QuestionType.SCALE)
    if not is_scale:
        try:
            vals = pd.to_numeric(qs.table["Варіант відповіді"], errors='coerce')
            if vals.notna().all() and vals.min() >= 0 and vals.max() <= 10:
                is_scale = True
        except: pass

    if is_scale:
        fig = plt.figure(figsize=(6.0, 4.5))
        bars = plt.bar(wrapped_labels, values, color='#4F81BD', width=BAR_WIDTH)
        plt.ylabel('Кількість')
        plt.grid(axis='y', linestyle='--', alpha=0.5)
        plt.xticks(rotation=0)
        for bar in bars:
            height = bar.get_height()
            plt.text(bar.get_x() + bar.get_width()/2., height + 0.1,
                     f'{int(height)}', ha='center', va='bottom', fontweight='bold')
    else:
        fig = plt.figure(figsize=(6.0, 5.0))
        colors = ['#4F81BD', '#C0504D', '#9BBB59', '#8064A2', '#4BACC6', '#F79646']
        c_arg = colors[:len(values)] if len(values) <= len(colors) else None
        
        wedges, texts, autotexts = plt.pie(
            values, labels=None, autopct='%1.1f%%', startangle=90,
            pctdistance=0.8, colors=c_arg, radius=1.1,
            textprops={'fontsize': FONT_SIZE_CHART}
        )
        for autotext in autotexts:
            autotext.set_color('white')
            autotext.set_weight('bold')
            import matplotlib.patheffects as path_effects
            autotext.set_path_effects([path_effects.withStroke(linewidth=2, foreground='#333333')])

        plt.axis('equal')
        cols = 2 if len(labels) > 2 else 1
        plt.legend(wrapped_labels, loc="upper center", bbox_to_anchor=(0.5, 0.0), ncol=cols, frameon=False, fontsize=10)

    plt.tight_layout()
    img_stream = io.BytesIO()
    plt.savefig(img_stream, format='png', dpi=CHART_DPI, bbox_inches='tight')
    plt.close(fig) 
    img_stream.seek(0)
    return img_stream

def build_pptx_report(original_df, sliced_df, summaries, range_info):
    prs = Presentation()
    # Слайд 1: Титул
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    if os.path.exists(LOGO_FILE):
        slide.shapes.add_picture(LOGO_FILE, Inches(0.2), Inches(0.2), width=Inches(1.2))

    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(0.4), Inches(8), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = UNIV_NAME
    p.font.bold = True
    p.font.size = Pt(16)
    p.alignment = PP_ALIGN.LEFT 

    title = slide.shapes.title
    title.text = "Звіт про результати опитування"
    subtitle = slide.placeholders[1]
    subtitle.text = f"Всього анкет: {len(original_df)}\nОброблено: {len(sliced_df)}\nДіапазон: {range_info}"

    # Слайди даних
    layout_index = 5 
    if len(prs.slide_layouts) <= 5: layout_index = len(prs.slide_layouts) - 1
    
    for qs in summaries:
        if qs.table.empty: continue
        
        slide = prs.slides.add_slide(prs.slide_layouts[layout_index])
        
        try:
            title = slide.shapes.title
            title.text = f"{qs.question.code}. {qs.question.text}"
            if len(title.text) > 60: title.text_frame.paragraphs[0].font.size = Pt(24)
            else: title.text_frame.paragraphs[0].font.size = Pt(32)
        except: pass

        # Таблиця
        rows = len(qs.table) + 1
        cols = 3
        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(2.0), Inches(4.5), Inches(0.8)).table
        set_table_grid_style(table)

        table.columns[0].width = Inches(2.5)
        table.columns[1].width = Inches(1.0)
        table.columns[2].width = Inches(1.0)

        # Header
        headers = ["Варіант", "Кільк.", "%"]
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = h
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(FONT_SIZE_HEADER)
            cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(220, 220, 220)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

        # Data
        for i, row in enumerate(qs.table.itertuples(index=False)):
            for j, val in enumerate(row):
                cell = table.cell(i+1, j)
                cell.text = str(val)
                cell.text_frame.paragraphs[0].font.size = Pt(FONT_SIZE_DATA)
                cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
                if j > 0: cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
                else: cell.text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)

        # Chart
        try:
            img_stream = create_chart_image(qs)
            slide.shapes.add_picture(img_stream, Inches(5.2), Inches(2.0), width=Inches(4.6))
        except: pass

 
    slide_layout = prs.slide_layouts[6] 
    slide = prs.slides.add_slide(slide_layout)

    if os.path.exists(LOGO_FILE):
        slide.shapes.add_picture(LOGO_FILE, Inches(4.25), Inches(1.0), width=Inches(1.5))

    tb_thanks = slide.shapes.add_textbox(Inches(0), Inches(3.0), Inches(10), Inches(1.5))
    tf_thanks = tb_thanks.text_frame
    p_thanks = tf_thanks.paragraphs[0]
    p_thanks.text = "Дякую за увагу!"
    p_thanks.alignment = PP_ALIGN.CENTER
    p_thanks.font.size = Pt(32)
    p_thanks.font.bold = True
    p_thanks.font.color.rgb = RGBColor(0, 0, 0) 

    tb_info = slide.shapes.add_textbox(Inches(0), Inches(4.5), Inches(10), Inches(2))
    tf_info = tb_info.text_frame
    
    lines = [
        "Створено за допомогою додатку студентки МПУіК – Каптар Діани.",
        "Керівник проєкту – доцент Фратавчан Валерій Григорович."
    ]

    for line in lines:
        p = tf_info.add_paragraph()
        p.text = line
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(18)  
        p.font.color.rgb = RGBColor(80, 80, 80) 
        p.space_after = Pt(10) 

    output = io.BytesIO()
    prs.save(output)
    return output.getvalue()